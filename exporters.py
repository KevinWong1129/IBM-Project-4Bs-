"""
exporters.py — Output Formatting Module
==========================================
Converts generated Markdown/plain-text answers into PDF, DOCX, or PPTX files.

Public API
----------
  from exporters import export
  export(answer_text, citations, prompt, output_path, fmt)
"""

from __future__ import annotations

import os
import textwrap
from datetime import datetime
from pathlib import Path


def export(answer: str, citations: str, prompt: str, output_path: str, fmt: str) -> None:
    fmt = fmt.lower()
    if fmt == "pdf":
        _export_pdf(answer, citations, prompt, output_path)
    elif fmt == "docx":
        _export_docx(answer, citations, prompt, output_path)
    elif fmt == "pptx":
        _export_pptx(answer, citations, prompt, output_path)
    else:
        raise ValueError(f"Unsupported format: {fmt!r}")


# ── PDF ───────────────────────────────────────────────────────────────────────

def _export_pdf(answer: str, citations: str, prompt: str, output_path: str) -> None:
    from fpdf import FPDF

    pdf = FPDF()
    pdf.set_margins(20, 20, 20)
    pdf.add_page()

    # Title
    pdf.set_font("Helvetica", "B", 16)
    pdf.cell(0, 10, "IBM Concert — Generated Brief", ln=True)
    pdf.set_font("Helvetica", "", 9)
    pdf.set_text_color(90, 90, 90)
    pdf.cell(0, 6, f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M')}  |  IBM Concert RAG System", ln=True)
    pdf.ln(4)

    # Prompt box
    pdf.set_font("Helvetica", "B", 10)
    pdf.set_text_color(31, 35, 40)
    pdf.cell(0, 8, "Request", ln=True)
    pdf.set_font("Helvetica", "I", 10)
    pdf.set_text_color(87, 96, 106)
    for line in textwrap.wrap(prompt, 90):
        pdf.cell(0, 6, line, ln=True)
    pdf.ln(4)

    # Divider
    pdf.set_draw_color(229, 231, 235)
    pdf.line(20, pdf.get_y(), 190, pdf.get_y())
    pdf.ln(4)

    # Body
    pdf.set_font("Helvetica", "B", 11)
    pdf.set_text_color(31, 35, 40)
    pdf.cell(0, 8, "Generated Content", ln=True)
    pdf.set_font("Helvetica", "", 10)
    for para in answer.split("\n"):
        para = para.strip()
        if not para:
            pdf.ln(3)
            continue
        # Detect headings (lines starting with #)
        if para.startswith("##"):
            pdf.set_font("Helvetica", "B", 11)
            pdf.cell(0, 7, para.lstrip("#").strip(), ln=True)
            pdf.set_font("Helvetica", "", 10)
        elif para.startswith("#"):
            pdf.set_font("Helvetica", "B", 12)
            pdf.cell(0, 8, para.lstrip("#").strip(), ln=True)
            pdf.set_font("Helvetica", "", 10)
        elif para.startswith("-") or para.startswith("•"):
            pdf.cell(5, 6, "•", ln=False)
            for chunk in textwrap.wrap(para.lstrip("-•").strip(), 85):
                pdf.cell(0, 6, "  " + chunk, ln=True)
        else:
            for chunk in textwrap.wrap(para, 90):
                pdf.cell(0, 6, chunk, ln=True)

    # Citations
    if citations:
        pdf.ln(6)
        pdf.line(20, pdf.get_y(), 190, pdf.get_y())
        pdf.ln(4)
        pdf.set_font("Helvetica", "B", 10)
        pdf.cell(0, 8, "Sources", ln=True)
        pdf.set_font("Helvetica", "", 9)
        pdf.set_text_color(87, 96, 106)
        for line in citations.split("\n"):
            pdf.cell(0, 5, line, ln=True)

    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    pdf.output(output_path)


# ── DOCX ──────────────────────────────────────────────────────────────────────

def _export_docx(answer: str, citations: str, prompt: str, output_path: str) -> None:
    from docx import Document
    from docx.shared import Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    doc = Document()

    # Title
    title = doc.add_heading("IBM Concert — Generated Brief", level=0)
    title.alignment = WD_ALIGN_PARAGRAPH.LEFT

    meta = doc.add_paragraph(f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M')}")
    meta.runs[0].font.size = Pt(9)
    meta.runs[0].font.color.rgb = RGBColor(87, 96, 106)

    doc.add_heading("Request", level=2)
    req = doc.add_paragraph(prompt)
    req.runs[0].italic = True

    doc.add_heading("Generated Content", level=2)

    for line in answer.split("\n"):
        line = line.strip()
        if not line:
            doc.add_paragraph("")
        elif line.startswith("## "):
            doc.add_heading(line[3:], level=2)
        elif line.startswith("# "):
            doc.add_heading(line[2:], level=1)
        elif line.startswith("- ") or line.startswith("• "):
            doc.add_paragraph(line.lstrip("-•").strip(), style="List Bullet")
        else:
            doc.add_paragraph(line)

    if citations:
        doc.add_heading("Sources", level=2)
        for line in citations.split("\n"):
            p = doc.add_paragraph(line)
            p.runs[0].font.size = Pt(9)

    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    doc.save(output_path)


# ── PPTX ──────────────────────────────────────────────────────────────────────

def _export_pptx(answer: str, citations: str, prompt: str, output_path: str) -> None:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]   # completely blank
    title_layout = prs.slide_layouts[0]   # title slide

    def add_title_slide(title_text: str, subtitle_text: str) -> None:
        slide   = prs.slides.add_slide(title_layout)
        title   = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = title_text
        subtitle.text = subtitle_text
        title.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(31, 35, 40)

    def add_content_slide(heading: str, bullets: list[str]) -> None:
        slide = prs.slides.add_slide(blank_layout)
        # Heading box
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.9))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = heading
        p.runs[0].font.bold  = True
        p.runs[0].font.size  = Pt(24)
        p.runs[0].font.color.rgb = RGBColor(31, 35, 40)
        # Bullets box
        bBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12), Inches(5.5))
        bTF  = bBox.text_frame
        bTF.word_wrap = True
        for i, bullet in enumerate(bullets):
            bp = bTF.paragraphs[0] if i == 0 else bTF.add_paragraph()
            bp.text = f"• {bullet}" if not bullet.startswith("•") else bullet
            bp.runs[0].font.size = Pt(18)
            bp.space_before = Pt(6)

    # Cover slide
    add_title_slide("IBM Concert — Generated Brief",
                    f"Request: {prompt[:120]}{'...' if len(prompt) > 120 else ''}")

    # Parse answer into slides: split on ## headings
    current_heading = "Overview"
    current_bullets: list[str] = []

    for line in answer.split("\n"):
        line = line.strip()
        if not line:
            continue
        if line.startswith("## ") or line.startswith("# "):
            if current_bullets:
                add_content_slide(current_heading, current_bullets)
            current_heading  = line.lstrip("#").strip()
            current_bullets  = []
        elif line.startswith("- ") or line.startswith("• "):
            current_bullets.append(line.lstrip("-•").strip())
        else:
            current_bullets.append(line)

    # Flush remaining content
    if current_bullets:
        add_content_slide(current_heading, current_bullets)

    # Sources slide
    if citations:
        add_content_slide("Sources", [c.strip() for c in citations.split("\n") if c.strip()])

    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
