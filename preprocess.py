"""
Pre-processing script: converts PDF and PPTX files in the Knowledge base
into clean Markdown files suitable for RAG ingestion.
"""

import os
import re
import pymupdf4llm
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN

KB_DIR = "Knowledge base"
OUT_DIR = "converted"  # output subfolder created inside KB_DIR

# DrawingML namespace used inside <p:graphicFrame> table XML
_NSMAP_A = "http://schemas.openxmlformats.org/drawingml/2006/main"


def pdf_to_markdown(pdf_path: str, out_path: str) -> None:
    """Convert a PDF to Markdown using pymupdf4llm."""
    print(f"  Converting PDF: {pdf_path}")
    md = pymupdf4llm.to_markdown(pdf_path)
    with open(out_path, "w", encoding="utf-8") as f:
        f.write(md)
    print(f"  → Written: {out_path}")


def table_to_markdown(shape) -> str:
    """
    Extract a <p:graphicFrame> table into a Markdown table string.

    WHY THIS IS NEEDED
    ──────────────────
    A PowerPoint table is stored as a <p:graphicFrame> element, not a <p:sp>
    (normal shape).  python-pptx reports has_text_frame = False for every
    graphicFrame, so the old shape loop skipped them with `continue` before
    even looking at their content.

    The actual text lives in the DrawingML subtree:
        <p:graphicFrame>
          <a:graphic>
            <a:graphicData>
              <a:tbl>          ← table root
                <a:tr>         ← row
                  <a:tc>       ← cell
                    <a:txBody>
                      <a:p><a:r><a:t>cell text</a:t>  ← text
    python-pptx exposes this via shape.table when shape.shape_type == TABLE,
    so we use that API instead of raw XML walking.

    NOTE: some exotic shape types (SmartArt, charts, embedded OLE objects)
    also live inside <p:graphicFrame> but raise NotImplementedError when
    .shape_type is accessed.  We guard with a try/except so those shapes are
    silently skipped rather than crashing the whole conversion.
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    try:
        if shape.shape_type != MSO_SHAPE_TYPE.TABLE:
            return ""
    except NotImplementedError:
        # Shape is a graphicFrame that isn't a table (e.g. chart, SmartArt,
        # embedded OLE object) — python-pptx can't identify its type, skip it.
        return ""

    tbl = shape.table
    if not tbl.rows:
        return ""

    md_rows = []
    for r, row in enumerate(tbl.rows):
        cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
        md_rows.append("| " + " | ".join(cells) + " |")
        # Insert separator after the first (header) row
        if r == 0:
            md_rows.append("| " + " | ".join(["---"] * len(cells)) + " |")

    return "\n".join(md_rows)


def shape_text(shape) -> str:
    """Extract formatted text from a PPTX shape's text frame."""
    if not shape.has_text_frame:
        return ""

    lines = []
    for para in shape.text_frame.paragraphs:
        text = para.text.strip()
        if not text:
            continue

        level = para.level  # 0 = top level, 1+ = indented bullets
        indent = "  " * level

        # Detect title-like paragraphs (large font or bold on level 0)
        # NOTE: font_size may be None when the size is inherited from the
        # theme/layout — guard against that before comparing to Pt(20).
        is_title = False
        if para.runs:
            run = para.runs[0]
            font_size = run.font.size
            is_bold = run.font.bold
            if font_size and font_size >= Pt(20) and level == 0:
                is_title = True
            elif is_bold and level == 0 and len(text) < 120:
                is_title = True

        if is_title:
            lines.append(f"### {text}")
        elif level == 0:
            lines.append(f"{indent}{text}")
        else:
            lines.append(f"{indent}- {text}")

    return "\n".join(lines)


def pptx_to_markdown(pptx_path: str, out_path: str) -> None:
    """Convert a PPTX to Markdown, one section per slide."""
    print(f"  Converting PPTX: {pptx_path}")
    prs = Presentation(pptx_path)
    sections = []

    for i, slide in enumerate(prs.slides, start=1):
        slide_lines = [f"## Slide {i}"]

        # Collect placeholder title first
        title_text = ""
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 0:  # title placeholder
                title_text = ph.text.strip()
                break

        if title_text:
            slide_lines.append(f"# {title_text}")

        # Process all shapes:
        #   • <p:sp>           → normal text box  → shape_text()
        #   • <p:graphicFrame> → table            → table_to_markdown()
        #                        (has_text_frame is always False for these,
        #                         so they were silently dropped before this fix)
        for shape in slide.shapes:
            # ── TABLE branch (graphicFrame) ──────────────────────────────────
            # Must be checked BEFORE the has_text_frame guard below, because
            # graphicFrame shapes always return has_text_frame = False and
            # would fall straight into `continue`.
            table_md = table_to_markdown(shape)
            if table_md:
                slide_lines.append(table_md)
                continue

            # ── NORMAL TEXT SHAPE branch ─────────────────────────────────────
            if not shape.has_text_frame:
                continue  # image, connector, group, etc. — nothing to extract
            # Skip the title placeholder we already handled
            try:
                if shape.placeholder_format and shape.placeholder_format.idx == 0:
                    continue
            except ValueError:
                pass  # shape.placeholder_format raises if not a placeholder
            text = shape_text(shape)
            if text:
                slide_lines.append(text)

        # Always keep the slide marker so numbering stays continuous
        # (previously empty slides like Slide 36 were dropped, breaking
        #  the sequence).  Add a note if the slide had no extractable text.
        content = "\n".join(slide_lines)
        if len(content.strip()) == len(f"## Slide {i}"):
            slide_lines.append("_[No extractable text — slide may be image-only]_")
            content = "\n".join(slide_lines)
        sections.append(content)

    md = "\n\n---\n\n".join(sections)
    with open(out_path, "w", encoding="utf-8") as f:
        f.write(md)
    print(f"  → Written: {out_path}")


def process_directory(root: str) -> None:
    """Walk the knowledge base and convert PDF/PPTX files to Markdown."""
    out_root = os.path.join(root, OUT_DIR)

    for dirpath, _, filenames in os.walk(root):
        # Don't recurse into the output folder itself
        if os.path.abspath(dirpath).startswith(os.path.abspath(out_root)):
            continue

        for filename in filenames:
            if filename.startswith("."):
                continue

            base, ext = os.path.splitext(filename)
            ext = ext.lower()
            src = os.path.join(dirpath, filename)

            # Mirror the subfolder structure inside out_root
            rel = os.path.relpath(dirpath, root)
            out_dir = os.path.join(out_root, rel) if rel != "." else out_root
            os.makedirs(out_dir, exist_ok=True)

            if ext == ".pdf":
                pdf_to_markdown(src, os.path.join(out_dir, base + "_pdf.md"))
            elif ext == ".pptx":
                pptx_to_markdown(src, os.path.join(out_dir, base + "_pptx.md"))


if __name__ == "__main__":
    print(f"Starting pre-processing under: {KB_DIR}\n")
    process_directory(KB_DIR)
    print("\nDone.")
